# -*- coding: utf-8 -*-
"""
Script tao slide PowerPoint v2 — Premium Modern Design
Palette: Deep navy + Purple gradient + Warm coral + Mint green
Chay: python generate_pptx.py
"""
import sys, io, os
if sys.platform == "win32":
    try:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')
    except Exception:
        pass

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ============================================================
# PALETTE V2 — Tinh te, hai hoa hon
# ============================================================
BG           = RGBColor(0x0F, 0x0E, 0x17)   # #0F0E17 - nen chinh (toi dam)
BG2          = RGBColor(0x16, 0x15, 0x24)   # #161524 - nen card
BG3          = RGBColor(0x1E, 0x1C, 0x31)   # #1E1C31 - nen card sang hon
PURPLE       = RGBColor(0x7F, 0x5A, 0xF0)   # #7F5AF0 - tim chu dao
PURPLE_L     = RGBColor(0xA7, 0x8B, 0xFA)   # #A78BFA - tim nhat
CORAL        = RGBColor(0xFF, 0x79, 0x79)   # #FF7979 - cam hong am
MINT         = RGBColor(0x2C, 0xB6, 0x7D)   # #2CB67D - xanh ngoc
MINT_L       = RGBColor(0x7E, 0xE8, 0xB8)   # #7EE8B8 - xanh ngoc sang
SKY          = RGBColor(0x56, 0xCC, 0xF2)   # #56CCF2 - xanh troi
GOLD         = RGBColor(0xF2, 0xC9, 0x4C)   # #F2C94C - vang am
WHITE        = RGBColor(0xFF, 0xFF, 0xFE)   # #FFFFFE
GRAY         = RGBColor(0x94, 0xA1, 0xB2)   # #94A1B2 - xam nhe
DIM          = RGBColor(0x72, 0x75, 0x7E)   # #72757E - xam dam
BORDER       = RGBColor(0x2D, 0x2B, 0x42)   # #2D2B42 - vien nhe

# ============================================================
# HELPERS
# ============================================================

def set_bg(slide, color=BG):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, border_c=None, bw=Pt(0), radius=True):
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if border_c:
        sh.line.fill.solid()
        sh.line.fill.fore_color.rgb = border_c
        sh.line.width = bw
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh

def txt(slide, l, t, w, h, text, sz=16, c=WHITE, bold=False, it=False,
        align=PP_ALIGN.LEFT, fn='Segoe UI Semibold'):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = c
    p.font.bold = bold
    p.font.italic = it
    p.font.name = fn
    p.alignment = align
    return tb

def mtxt(slide, l, t, w, h, lines, ls=1.2):
    """lines: list of dict {text, sz, c, bold, it, fn, align, spa}"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, d in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = d.get('text', '')
        p.font.size = Pt(d.get('sz', 14))
        p.font.color.rgb = d.get('c', WHITE)
        p.font.bold = d.get('bold', False)
        p.font.italic = d.get('it', False)
        p.font.name = d.get('fn', 'Segoe UI')
        p.alignment = d.get('align', PP_ALIGN.LEFT)
        p.space_after = Pt(d.get('spa', 4))
        p.line_spacing = ls
    return tb

def accent_bar(slide, l, t, w, c=PURPLE, h=Pt(3)):
    return rect(slide, l, t, w, h, c, radius=False)

def page_num(slide, n, total):
    txt(slide, Inches(9.0), Inches(7.05), Inches(0.8), Inches(0.3),
        f"{n}/{total}", sz=9, c=DIM, align=PP_ALIGN.RIGHT, fn='Segoe UI')

def heading(slide, title, sub=""):
    txt(slide, Inches(0.7), Inches(0.35), Inches(8.6), Inches(0.55),
        title, sz=26, c=WHITE, bold=True)
    accent_bar(slide, Inches(0.7), Inches(0.88), Inches(0.9), PURPLE)
    if sub:
        txt(slide, Inches(0.7), Inches(1.0), Inches(8.6), Inches(0.35),
            sub, sz=12, c=GRAY, it=True, fn='Segoe UI')

def stat(slide, l, t, w, h, val, lab, ac=PURPLE):
    rect(slide, l, t, w, h, BG2, border_c=ac, bw=Pt(1.5))
    txt(slide, l, t + Inches(0.12), w, Inches(0.45),
        val, sz=26, c=ac, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, l, t + Inches(0.58), w, Inches(0.35),
        lab, sz=10, c=GRAY, align=PP_ALIGN.CENTER, fn='Segoe UI')

def card(slide, l, t, w, h, title, body, ac=PURPLE, tsz=14, bsz=12):
    rect(slide, l, t, w, h, BG2, border_c=BORDER, bw=Pt(1))
    rect(slide, l, t, Pt(4), h, ac, radius=False)
    txt(slide, l + Inches(0.2), t + Inches(0.08), w - Inches(0.3), Inches(0.35),
        title, sz=tsz, c=ac, bold=True)
    y = t + Inches(0.42)
    for line in body:
        txt(slide, l + Inches(0.2), y, w - Inches(0.3), Inches(0.25),
            line, sz=bsz, c=GRAY, fn='Segoe UI')
        y += Inches(0.26)

def tbl(slide, l, t, w, h, headers, rows, cw=None):
    ts = slide.shapes.add_table(1 + len(rows), len(headers), l, t, w, h)
    table = ts.table
    if cw:
        for j, cwidth in enumerate(cw):
            table.columns[j].width = Inches(cwidth)
    for j, hdr in enumerate(headers):
        c = table.cell(0, j)
        c.text = hdr
        c.fill.solid()
        c.fill.fore_color.rgb = PURPLE
        p = c.text_frame.paragraphs[0]
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = 'Segoe UI Semibold'
        p.alignment = PP_ALIGN.CENTER
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, rd in enumerate(rows):
        for j, v in enumerate(rd):
            c = table.cell(i + 1, j)
            c.text = str(v)
            c.fill.solid()
            c.fill.fore_color.rgb = BG3 if i % 2 == 0 else BG2
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = WHITE
            p.font.name = 'Segoe UI'
            p.alignment = PP_ALIGN.CENTER
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
    return ts

def deco_circle(slide, l, t, sz, c, op=0.15):
    """Vong tron trang tri nen — dung lam diem nhan tham my."""
    from lxml import etree
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    sh.fill.solid()
    sh.fill.fore_color.rgb = c
    # Set transparency via XML
    spPr = sh._element.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            alpha = etree.SubElement(srgb, qn('a:alpha'))
            alpha.set('val', str(int(op * 100000)))
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


# ============================================================
# BUILD
# ============================================================
def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    T = 14  # total slides

    # ──────────── SLIDE 1: TITLE ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    # Decorative circles
    deco_circle(s, Inches(7.5), Inches(-1), Inches(4), PURPLE, 0.06)
    deco_circle(s, Inches(8.5), Inches(4.5), Inches(3.5), MINT, 0.04)
    deco_circle(s, Inches(-1.5), Inches(5), Inches(3), CORAL, 0.04)
    # Top bar
    accent_bar(s, Inches(0), Inches(0), Inches(10), PURPLE, Pt(4))
    # Content
    mtxt(s, Inches(0.9), Inches(1.6), Inches(8), Inches(4.5), [
        {'text': 'BLOOM FILTER', 'sz': 48, 'c': WHITE, 'bold': True, 'fn': 'Segoe UI Semibold', 'spa': 0},
        {'text': 'SEMI-JOIN OPTIMIZER', 'sz': 48, 'c': PURPLE_L, 'bold': True, 'fn': 'Segoe UI Semibold', 'spa': 20},
        {'text': 'Tối ưu hóa Phép Nối Phân Tán cho Hệ thống', 'sz': 16, 'c': GRAY, 'spa': 2},
        {'text': 'Subscribers & WebLogs', 'sz': 16, 'c': GRAY, 'spa': 30},
        {'text': 'Đề tài #15  ·  Category 3: Tối ưu hóa Truy vấn Phân tán', 'sz': 14, 'c': GOLD, 'bold': True, 'spa': 6},
        {'text': 'Môn: Cơ Sở Dữ Liệu Phân Tán', 'sz': 13, 'c': GRAY, 'spa': 10},
        {'text': '[Điền tên nhóm  ·  Thành viên 1  ·  Thành viên 2]', 'sz': 13, 'c': DIM, 'it': True},
    ], ls=1.15)
    page_num(s, 1, T)

    # ──────────── SLIDE 2: AGENDA ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(8), Inches(0.5), Inches(3), PURPLE, 0.04)
    heading(s, "NỘI DUNG TRÌNH BÀY")

    items = [
        ("01", "Bài toán & Mục tiêu", "Vấn đề JOIN phân tán, lãng phí 80% băng thông"),
        ("02", "Kiến trúc hệ thống", "2 Sites phân tán giao tiếp qua REST API"),
        ("03", "Bloom Filter", "Cấu trúc dữ liệu xác suất — cốt lõi giải thuật"),
        ("04", "Quy trình Semi-Join 5 bước", "Pipeline tối ưu truyền tải dữ liệu"),
        ("05", "Kết quả thực nghiệm", "So sánh Naive Join vs BF Semi-Join"),
        ("06", "Phân tích lý thuyết 1M/10M", "Tính toán theo Özsu & Valduriez"),
        ("07", "Khả năng chịu lỗi", "Demo Kill Node B & Recovery"),
        ("08", "Kết luận & Hướng phát triển", "Tổng kết kết quả đồ án"),
    ]
    colors = [PURPLE_L, SKY, MINT_L, GOLD, CORAL, PURPLE_L, SKY, MINT_L]
    y = Inches(1.55)
    for i, (num, title, desc) in enumerate(items):
        ac = colors[i % len(colors)]
        # Number badge
        badge = rect(s, Inches(0.7), y, Inches(0.42), Inches(0.42), ac)
        tf = badge.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BG
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = 'Segoe UI Semibold'
        # Text
        txt(s, Inches(1.3), y - Inches(0.02), Inches(6), Inches(0.3),
            title, sz=15, c=WHITE, bold=True)
        txt(s, Inches(1.3), y + Inches(0.26), Inches(7), Inches(0.25),
            desc, sz=11, c=DIM, fn='Segoe UI')
        y += Inches(0.68)
    page_num(s, 2, T)

    # ──────────── SLIDE 3: BÀI TOÁN ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(-0.5), Inches(5), Inches(2.5), CORAL, 0.05)
    heading(s, "BÀI TOÁN & MỤC TIÊU", "Vì sao cần tối ưu JOIN phân tán?")

    # Query
    rect(s, Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.65), BG3, border_c=PURPLE, bw=Pt(1.5))
    txt(s, Inches(0.85), Inches(1.46), Inches(8.3), Inches(0.5),
        'SELECT S.*, W.*  FROM Subscribers S  JOIN  WebLogs W  ON  S.user_id = W.user_id',
        sz=12, c=GOLD, fn='Cascadia Mono')

    # Problem
    mtxt(s, Inches(0.6), Inches(2.3), Inches(8.8), Inches(2.4), [
        {'text': '⚠  VẤN ĐỀ', 'sz': 17, 'c': CORAL, 'bold': True, 'spa': 10},
        {'text': '•  Site A giữ bảng Subscribers — 1,000,000 thuê bao trả phí', 'sz': 13, 'c': WHITE, 'spa': 3},
        {'text': '•  Site B giữ bảng WebLogs — 10,000,000 nhật ký truy cập web', 'sz': 13, 'c': WHITE, 'spa': 3},
        {'text': '•  80% lượt truy cập là khách vãng lai (guest) → không cần JOIN', 'sz': 13, 'c': WHITE, 'spa': 8},
        {'text': '→  Naive Join: Gửi toàn bộ 10M rows (~1,430 MB) qua mạng', 'sz': 13, 'c': CORAL, 'spa': 3},
        {'text': '→  Trong đó 80% (~1,144 MB) là dữ liệu thừa — hoàn toàn lãng phí!', 'sz': 13, 'c': CORAL, 'bold': True, 'spa': 14},
    ], ls=1.15)

    # Solution
    rect(s, Inches(0.6), Inches(4.95), Inches(8.8), Inches(1.1), BG2, border_c=MINT, bw=Pt(1.5))
    mtxt(s, Inches(0.85), Inches(5.0), Inches(8.4), Inches(1.0), [
        {'text': '✅  GIẢI PHÁP: Bloom Filter Semi-Join', 'sz': 17, 'c': MINT_L, 'bold': True, 'spa': 6},
        {'text': 'Lọc dữ liệu thừa ngay tại Site B → chỉ gửi ~20% cần thiết (~307 MB) → tiết kiệm ~78.6%', 'sz': 13, 'c': WHITE},
    ], ls=1.2)
    page_num(s, 3, T)

    # ──────────── SLIDE 4: KIẾN TRÚC ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "KIẾN TRÚC HỆ THỐNG PHÂN TÁN", "2 Sites độc lập — giao tiếp qua HTTP REST API")

    # Site A
    rect(s, Inches(0.4), Inches(1.8), Inches(3.6), Inches(3.5), BG2, border_c=PURPLE, bw=Pt(2))
    txt(s, Inches(0.4), Inches(1.4), Inches(3.6), Inches(0.35),
        "SITE A — Trụ sở chính (port 5000)", sz=13, c=PURPLE_L, bold=True, align=PP_ALIGN.CENTER)
    mtxt(s, Inches(0.7), Inches(2.05), Inches(3), Inches(2.8), [
        {'text': '📋  Subscribers', 'sz': 15, 'c': WHITE, 'bold': True, 'spa': 6},
        {'text': '1,000,000 thuê bao trả phí', 'sz': 11, 'c': GRAY, 'spa': 3},
        {'text': 'user_id · full_name · plan · fee', 'sz': 10, 'c': DIM, 'fn': 'Cascadia Mono', 'spa': 14},
        {'text': '🔧  Vai trò', 'sz': 13, 'c': WHITE, 'bold': True, 'spa': 5},
        {'text': '•  Tạo Bloom Filter từ user_ids', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Gửi BF → Site B qua HTTP', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Nhận filtered logs về', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Inner Join cuối → kết quả 100%', 'sz': 11, 'c': GRAY},
    ], ls=1.1)

    # Arrows
    txt(s, Inches(4.1), Inches(2.5), Inches(1.8), Inches(0.35),
        "① BF bit-vector →", sz=10, c=GOLD, bold=True, align=PP_ALIGN.CENTER, fn='Cascadia Mono')
    ar1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.15), Inches(2.95), Inches(1.7), Inches(0.22))
    ar1.fill.solid(); ar1.fill.fore_color.rgb = GOLD; ar1.line.fill.background(); ar1.shadow.inherit = False

    txt(s, Inches(4.1), Inches(3.7), Inches(1.8), Inches(0.35),
        "② ← Filtered Logs", sz=10, c=MINT_L, bold=True, align=PP_ALIGN.CENTER, fn='Cascadia Mono')
    ar2 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.15), Inches(4.1), Inches(1.7), Inches(0.22))
    ar2.fill.solid(); ar2.fill.fore_color.rgb = MINT; ar2.line.fill.background(); ar2.shadow.inherit = False

    # Site B
    rect(s, Inches(6.0), Inches(1.8), Inches(3.6), Inches(3.5), BG2, border_c=MINT, bw=Pt(2))
    txt(s, Inches(6.0), Inches(1.4), Inches(3.6), Inches(0.35),
        "SITE B — Web Server (port 5001)", sz=13, c=MINT_L, bold=True, align=PP_ALIGN.CENTER)
    mtxt(s, Inches(6.3), Inches(2.05), Inches(3), Inches(2.8), [
        {'text': '📋  WebLogs', 'sz': 15, 'c': WHITE, 'bold': True, 'spa': 6},
        {'text': '10,000,000 nhật ký truy cập', 'sz': 11, 'c': GRAY, 'spa': 3},
        {'text': 'log_id · user_id · page · status', 'sz': 10, 'c': DIM, 'fn': 'Cascadia Mono', 'spa': 14},
        {'text': '🔧  Vai trò', 'sz': 13, 'c': WHITE, 'bold': True, 'spa': 5},
        {'text': '•  Nhận BF từ Site A', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Lọc WebLogs qua BF', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Loại bỏ ~80% guest logs', 'sz': 11, 'c': GRAY, 'spa': 2},
        {'text': '•  Gửi filtered logs về A', 'sz': 11, 'c': GRAY},
    ], ls=1.1)

    # Tech stack bar
    rect(s, Inches(0.4), Inches(5.7), Inches(9.2), Inches(0.5), BG3, border_c=BORDER, bw=Pt(1))
    txt(s, Inches(0.4), Inches(5.75), Inches(9.2), Inches(0.4),
        "Python 3.8+  ·  Flask REST  ·  MurmurHash3  ·  bitarray  ·  Pandas  ·  Chart.js",
        sz=11, c=DIM, align=PP_ALIGN.CENTER, fn='Segoe UI')
    page_num(s, 4, T)

    # ──────────── SLIDE 5: BLOOM FILTER ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(8), Inches(5), Inches(3), PURPLE, 0.04)
    heading(s, "BLOOM FILTER", "Cấu trúc dữ liệu xác suất — cốt lõi của giải pháp")

    txt(s, Inches(0.7), Inches(1.35), Inches(9), Inches(0.35),
        "Bit Array (m bits)  +  k hàm băm (Double Hashing — MurmurHash3)",
        sz=14, c=WHITE, bold=True)

    # Bit array
    bits = [0,0,1,0,1,0,0,1,0,1,1,0,0,1,0,0,1,0,0,1,0,1,0,0]
    x0 = Inches(0.7)
    yb = Inches(1.85)
    cw_bit = Inches(0.36)
    ch_bit = Inches(0.36)
    for i, b in enumerate(bits):
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x0 + cw_bit * i, yb, cw_bit, ch_bit)
        r.fill.solid()
        r.fill.fore_color.rgb = PURPLE if b else RGBColor(0x1A, 0x18, 0x2E)
        r.line.fill.solid()
        r.line.fill.fore_color.rgb = RGBColor(0x33, 0x30, 0x50)
        r.line.width = Pt(0.75)
        r.shadow.inherit = False
        tf = r.text_frame
        tf.paragraphs[0].text = str(b)
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE if b else DIM
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = 'Cascadia Mono'

    # Formulas
    mtxt(s, Inches(0.7), Inches(2.5), Inches(9), Inches(1.5), [
        {'text': 'Công thức tối ưu (Özsu & Valduriez):', 'sz': 14, 'c': GOLD, 'bold': True, 'spa': 8},
        {'text': '   m = −(n × ln(p)) / (ln2)²          →  Kích thước bit array tối ưu', 'sz': 12, 'c': WHITE, 'fn': 'Cascadia Mono', 'spa': 4},
        {'text': '   k = (m/n) × ln(2)                   →  Số hàm băm tối ưu', 'sz': 12, 'c': WHITE, 'fn': 'Cascadia Mono', 'spa': 4},
        {'text': '   FPR = (1 − e^(−kn/m))^k             →  Tỷ lệ False Positive thực tế', 'sz': 12, 'c': WHITE, 'fn': 'Cascadia Mono'},
    ], ls=1.2)

    # Cards
    card(s, Inches(0.5), Inches(4.7), Inches(4.3), Inches(1.5),
         "✅  False Negative = 0%",
         ["Bloom Filter KHÔNG BAO GIỜ bỏ sót phần tử",
          "thực sự có trong tập hợp. An toàn tuyệt đối."],
         ac=MINT, bsz=12)
    card(s, Inches(5.2), Inches(4.7), Inches(4.3), Inches(1.5),
         "⚠  False Positive > 0%  (kiểm soát được)",
         ["Có thể nhận nhầm → nhưng bị loại bỏ hoàn toàn",
          "ở bước Inner Join cuối cùng. Kết quả vẫn 100%."],
         ac=GOLD, bsz=12)
    page_num(s, 5, T)

    # ──────────── SLIDE 6: PIPELINE 5 BƯỚC ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "QUY TRÌNH SEMI-JOIN — 5 BƯỚC", "Pipeline tối ưu truyền tải dữ liệu giữa 2 Sites")

    steps = [
        ("1", "Tạo BF\ntại Site A",  PURPLE_L, "n=1M\nFPR=1%"),
        ("2", "Gửi BF\nA → B",       GOLD,     "~1.14\nMB"),
        ("3", "Lọc tại\nSite B",      MINT_L,   "Loại\n~80%"),
        ("4", "Gửi kết quả\nB → A",   GOLD,     "~306\nMB"),
        ("5", "Inner Join\ntại Site A", PURPLE_L, "100%\nchính xác"),
    ]
    x0 = Inches(0.25)
    sw = Inches(1.65)
    gap = Inches(0.3)
    ys = Inches(1.6)

    for i, (num, lab, clr, det) in enumerate(steps):
        x = x0 + i * (sw + gap)
        rect(s, x, ys, sw, Inches(1.8), BG2, border_c=clr, bw=Pt(2))
        # Circle
        ci = s.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.57), ys - Inches(0.22), Inches(0.5), Inches(0.5))
        ci.fill.solid(); ci.fill.fore_color.rgb = clr
        ci.line.fill.background(); ci.shadow.inherit = False
        tf = ci.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BG
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].font.name = 'Segoe UI Semibold'
        # Label
        txt(s, x + Inches(0.08), ys + Inches(0.4), sw - Inches(0.16), Inches(0.65),
            lab, sz=12, c=WHITE, bold=True, align=PP_ALIGN.CENTER)
        # Detail
        txt(s, x + Inches(0.05), ys + Inches(1.2), sw - Inches(0.1), Inches(0.45),
            det, sz=10, c=clr, align=PP_ALIGN.CENTER, fn='Segoe UI')
        # Arrow
        if i < 4:
            ar = s.shapes.add_shape(MSO_SHAPE.CHEVRON, x + sw + Inches(0.03), ys + Inches(0.7),
                                    Inches(0.24), Inches(0.3))
            ar.fill.solid(); ar.fill.fore_color.rgb = DIM
            ar.line.fill.background(); ar.shadow.inherit = False

    # Summary bar
    rect(s, Inches(0.4), Inches(3.8), Inches(9.2), Inches(0.55), BG2, border_c=MINT, bw=Pt(1.5))
    txt(s, Inches(0.4), Inches(3.86), Inches(9.2), Inches(0.4),
        "~307 MB thay vì ~1,430 MB  →  Tiết kiệm ~78.5% băng thông  ·  Kết quả chính xác 100%",
        sz=14, c=MINT_L, bold=True, align=PP_ALIGN.CENTER)

    # Table
    tbl(s, Inches(0.4), Inches(4.7), Inches(9.2), Inches(2),
        ["Bước", "Vị trí", "Hành động", "Chi phí mạng"],
        [
            ["1", "Site A", "Tạo BF từ user_ids (n=1M, FPR=1%)", "0"],
            ["2", "A → B", "Gửi BF bit-vector qua HTTP", "~1.14 MB"],
            ["3", "Site B", "Lọc WebLogs qua BF, loại ~80% guest", "0"],
            ["4", "B → A", "Gửi filtered logs (~2M + FP rows)", "~306 MB"],
            ["5", "Site A", "Inner Join cuối → loại FP, kết quả 100%", "0"],
        ], cw=[0.6, 1, 4.6, 1.5])
    page_num(s, 6, T)

    # ──────────── SLIDE 7: KẾT QUẢ ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "KẾT QUẢ THỰC NGHIỆM", "100K Subscribers + 1M WebLogs  ·  Overlap 20%")

    stat(s, Inches(0.4), Inches(1.35), Inches(2.1), Inches(1.0), "~80%", "Bandwidth Saved", MINT)
    stat(s, Inches(2.7), Inches(1.35), Inches(2.1), Inches(1.0), "666×", "Savings Leverage", PURPLE_L)
    stat(s, Inches(5.0), Inches(1.35), Inches(2.1), Inches(1.0), "100%", "Độ chính xác", GOLD)
    stat(s, Inches(7.3), Inches(1.35), Inches(2.3), Inches(1.0), "0%", "False Negative", MINT)

    tbl(s, Inches(0.4), Inches(2.7), Inches(9.2), Inches(2),
        ["Chiến lược", "Dòng gửi", "FP", "Mạng (KB)", "Tiết kiệm"],
        [
            ["Naive Join", "1,000,000", "0", "146,484", "baseline"],
            ["BF Semi-Join (FPR 10%)", "~280,000", "~80,000", "~41,016", "~72%"],
            ["BF Semi-Join (FPR 5%)", "~240,000", "~40,000", "~35,156", "~76%"],
            ["BF Semi-Join (FPR 1%)", "~208,000", "~8,000", "~30,469", "~79%"],
            ["BF Semi-Join (FPR 0.1%)", "~200,800", "~800", "~29,414", "~80%"],
        ], cw=[2.5, 1.5, 1.2, 1.5, 1.2])

    txt(s, Inches(0.5), Inches(5.0), Inches(9), Inches(0.35),
        "Bảng 2 — METRIC CHÍNH: Bytes Saved vs BF Size (m bits)", sz=14, c=WHITE, bold=True)
    tbl(s, Inches(0.4), Inches(5.4), Inches(9.2), Inches(1.5),
        ["FPR", "BF size (m bits)", "BF (KB)", "Bytes Saved (KB)", "Leverage"],
        [
            ["10%", "479,253", "58.5", "~105,468", "~1,803×"],
            ["5%", "623,527", "76.1", "~111,328", "~1,463×"],
            ["1%", "958,506", "117.0", "~116,015", "~991×"],
            ["0.1%", "1,437,759", "175.5", "~117,070", "~667×"],
        ], cw=[1.2, 2.2, 1.5, 2, 1.5])
    page_num(s, 7, T)

    # ──────────── SLIDE 8: METRIC ANALYSIS ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(8.5), Inches(5.5), Inches(2.5), GOLD, 0.04)
    heading(s, "PHÂN TÍCH METRIC CHÍNH", "Bytes Saved vs. Size of Bit-Vector (m bits)")

    card(s, Inches(0.4), Inches(1.4), Inches(4.4), Inches(2.5),
         "📊  Nhận xét quan trọng",
         ["",
          "•  Leverage cực cao: 1 byte BF → tiết kiệm",
          "   667 — 1,803 bytes bandwidth mạng.",
          "",
          "•  FPR thấp hơn → BF lớn hơn → ít FP hơn",
          "   → tiết kiệm nhiều hơn, nhưng leverage giảm.",
          "",
          "•  FPR = 1%: tốn thêm chỉ ~0.8% bandwidth",
          "   do False Positive — hoàn toàn chấp nhận được."],
         ac=SKY, bsz=12)

    card(s, Inches(5.2), Inches(1.4), Inches(4.4), Inches(2.5),
         "🎯  Trade-off tối ưu",
         ["",
          "•  FPR = 1% là điểm cân bằng tốt nhất:",
          "   BF = 117 KB → Saved = 116,015 KB",
          "   → Leverage gần 1,000 lần!",
          "",
          "•  So với Semi-Join truyền thống (gửi IDs):",
          "   307 MB < 317 MB — BF nhẹ hơn 10× so với",
          "   danh sách 1M user_ids (10 MB vs 1.14 MB)."],
         ac=GOLD, bsz=12)

    # FPR Impact table
    txt(s, Inches(0.5), Inches(4.2), Inches(9), Inches(0.35),
        "FPR Impact on Wasted Bandwidth:", sz=14, c=WHITE, bold=True)
    tbl(s, Inches(0.4), Inches(4.6), Inches(9.2), Inches(1.5),
        ["FPR Target", "FP Rows", "FP Bytes (KB)", "% Bandwidth thêm do FP"],
        [
            ["10%", "~80,000", "~11,719", "~8.0%"],
            ["5%", "~40,000", "~5,859", "~4.0%"],
            ["1%", "~8,000", "~1,172", "~0.8%"],
            ["0.1%", "~800", "~117", "~0.08%"],
        ], cw=[2, 2, 2.5, 3])
    page_num(s, 8, T)

    # ──────────── SLIDE 9: LÝ THUYẾT 1M/10M ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "PHÂN TÍCH LÝ THUYẾT — QUY MÔ 1M / 10M", "Tính toán theo công thức BF cho quy mô chính xác đề bài")

    tbl(s, Inches(0.3), Inches(1.4), Inches(9.4), Inches(2),
        ["FPR", "m (bits)", "BF (MB)", "Total (MB)", "Saved (MB)", "Saved %", "Leverage"],
        [
            ["10%", "4,792,530", "0.57", "401.97", "1,028.03", "71.9%", "1,803×"],
            ["5%", "6,235,224", "0.74", "344.83", "1,085.17", "75.9%", "1,462×"],
            ["1%", "9,585,059", "1.14", "306.54", "1,123.46", "78.6%", "983×"],
            ["0.5%", "11,035,332", "1.31", "297.40", "1,132.60", "79.2%", "862×"],
            ["0.1%", "14,377,588", "1.71", "288.26", "1,141.74", "79.8%", "667×"],
        ], cw=[0.8, 1.6, 1.0, 1.3, 1.3, 1.0, 1.0])

    stat(s, Inches(0.3), Inches(3.8), Inches(2.9), Inches(1.0), "~1,430 MB", "Naive Join (gửi hết)", CORAL)
    stat(s, Inches(3.5), Inches(3.8), Inches(2.9), Inches(1.0), "~307 MB", "BF Semi-Join (FPR=1%)", MINT)
    stat(s, Inches(6.7), Inches(3.8), Inches(2.9), Inches(1.0), "~1,123 MB", "Tiết kiệm 78.6%", GOLD)

    card(s, Inches(0.3), Inches(5.1), Inches(9.4), Inches(1.7),
         "📖  Đánh giá theo Özsu & Valduriez (Chương 7 — Distributed Query Processing)",
         ["•  BF Semi-Join tối ưu hơn Semi-Join truyền thống: 307 MB vs 317 MB — giảm thêm 10 MB.",
          "•  Chi phí truyền tải: C_total = C₀ + C₁ × |data|  →  giảm |data| = giảm chi phí.",
          "•  Trade-off: Tính hash local (rẻ) để giảm data truyền mạng (đắt) — đúng nguyên lý tối ưu.",
          "•  Kết quả tương đương centralized JOIN — tính đúng đắn (Correctness) = 100%."],
         ac=PURPLE_L, bsz=12)
    page_num(s, 9, T)

    # ──────────── SLIDE 10: SO SÁNH ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "SO SÁNH CÁC CHIẾN LƯỢC JOIN", "Theo Özsu & Valduriez, Chương 7")

    tbl(s, Inches(0.3), Inches(1.4), Inches(9.4), Inches(1.4),
        ["Chiến lược", "Chi phí mạng", "Ưu điểm", "Nhược điểm"],
        [
            ["Ship Whole (Naive)", "~1,430 MB", "Đơn giản", "Lãng phí cực lớn"],
            ["Semi-Join (gửi IDs)", "~317 MB", "Giảm ~78%", "Vẫn tốn 10 MB gửi IDs"],
            ["BF Semi-Join (đề án)", "~307 MB", "Giảm 78.6%, BF nhỏ gọn", "FPR > 0 (= 0 ở kết quả)"],
        ], cw=[2.5, 1.8, 2.5, 2.8])

    txt(s, Inches(0.5), Inches(3.1), Inches(9), Inches(0.35),
        "Chi phí truyền tải:  C_total = C₀ + C₁ × |data|", sz=14, c=GOLD, bold=True, fn='Cascadia Mono')

    tbl(s, Inches(0.3), Inches(3.6), Inches(9.4), Inches(1.2),
        ["Phương pháp", "|data| truyền", "Số lần truyền"],
        [
            ["Naive Join", "10M × 150 bytes = 1,430 MB", "1 chiều (B→A)"],
            ["Semi-Join", "1M×10B + 2M×150B ≈ 310 MB", "2 chiều (A→B→A)"],
            ["BF Semi-Join", "1.14 MB + 2M×150B ≈ 307 MB", "2 chiều, chiều đi cực nhỏ"],
        ], cw=[2.2, 4, 3])

    card(s, Inches(0.3), Inches(5.1), Inches(9.4), Inches(1.6),
         "✅  Tính đúng đắn (Correctness)",
         ["•  False Negative = 0%  →  KHÔNG BAO GIỜ bỏ sót subscriber  →  mọi matching tuple đều được gửi về.",
          "•  False Positive bị loại bỏ  →  Inner Join cuối tại Site A loại tự động.",
          "•  Kết quả = 100% chính xác  →  Tương đương centralized JOIN."],
         ac=MINT, bsz=12)
    page_num(s, 10, T)

    # ──────────── SLIDE 11: FAULT TOLERANCE ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(-1), Inches(5), Inches(3), CORAL, 0.04)
    heading(s, "KHẢ NĂNG CHỊU LỖI", "Failure Scenario: Kill Node B")

    card(s, Inches(0.3), Inches(1.35), Inches(4.5), Inches(2.6),
         "✅  Luồng bình thường (Online)",
         ["", "①  Site A tạo Bloom Filter",
          "②  Gửi BF → Site B (HTTP POST)",
          "③  Site B lọc WebLogs qua BF",
          "④  Gửi filtered logs về Site A",
          "⑤  Site A Inner Join → 100%",
          "", "→  Kết quả chính xác, hoàn tất thành công."],
         ac=MINT, bsz=12)

    card(s, Inches(5.2), Inches(1.35), Inches(4.5), Inches(2.6),
         "❌  Node B bị KILL (Offline)",
         ["", "①  Site A tạo BF thành công",
          "②  Gửi BF → Site B...",
          "✖  Connection Timeout (HTTP 504)!",
          "③④⑤  ABORTED — không thể tiếp tục",
          "",
          "→  Hủy giao dịch + Rollback trạng thái",
          "→  Không treo kết nối, không mất dữ liệu."],
         ac=CORAL, bsz=12)

    card(s, Inches(0.3), Inches(4.2), Inches(9.4), Inches(2.6),
         "🔄  Quy trình xử lý lỗi (Fault Tolerance Design)",
         ["",
          "1.  Phát hiện lỗi (Fault Detection):  Site A thiết lập Timeout. Quá thời gian → xác định Node B sập.",
          "2.  Hủy bỏ & Khôi phục (Abort & Rollback):  Dừng JOIN, giải phóng tài nguyên, không cập nhật kết quả.",
          "3.  Cảnh báo trực quan:  Web Dashboard hiển thị pipeline đỏ + banner cảnh báo sập giao dịch.",
          "4.  Khôi phục (Recovery):  Nhấn \"Khôi phục Node B\" → kết nối lại → chạy lại bình thường.",
          "",
          "→  Đảm bảo tính nhất quán hệ thống phân tán — không treo kết nối, không gây mất dữ liệu!"],
         ac=SKY, bsz=12)
    page_num(s, 11, T)

    # ──────────── SLIDE 12: DASHBOARD ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    heading(s, "WEB DASHBOARD", "Giao diện trực quan tương tác — http://localhost:5000")

    feats = [
        ("🎛️  Cấu hình động", "Slider điều chỉnh Subscribers,\nWebLogs, Overlap % trực tiếp\ntrên trình duyệt.", PURPLE_L),
        ("📊  Biểu đồ real-time", "Chart.js: So sánh bandwidth,\nFPR, Leverage, Break-even\nPoint với animation.", SKY),
        ("🔀  Pipeline trực quan", "Animation 5 bước Semi-Join.\nHiển thị lỗi đỏ khi Node B\nbị kill + Recovery.", GOLD),
        ("🌐  Phân tán thực tế", "2 Flask servers độc lập\n(port 5000 + 5001) giao tiếp\nqua HTTP REST API.", MINT_L),
    ]
    for i, (title, desc, clr) in enumerate(feats):
        col, row = i % 2, i // 2
        x = Inches(0.4) + col * Inches(4.8)
        y = Inches(1.4) + row * Inches(2.5)
        rect(s, x, y, Inches(4.4), Inches(2.1), BG2, border_c=clr, bw=Pt(1.5))
        txt(s, x + Inches(0.2), y + Inches(0.15), Inches(4), Inches(0.35),
            title, sz=15, c=clr, bold=True)
        txt(s, x + Inches(0.2), y + Inches(0.55), Inches(4), Inches(1.3),
            desc, sz=12, c=GRAY, fn='Segoe UI')

    rect(s, Inches(0.4), Inches(6.6), Inches(9.2), Inches(0.45), BG3, border_c=BORDER, bw=Pt(1))
    txt(s, Inches(0.4), Inches(6.64), Inches(9.2), Inches(0.35),
        "Terminal 1:  python site_b_server.py   ·   Terminal 2:  python app.py   ·   Browser:  localhost:5000",
        sz=11, c=DIM, align=PP_ALIGN.CENTER, fn='Cascadia Mono')
    page_num(s, 12, T)

    # ──────────── SLIDE 13: KẾT LUẬN ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(8), Inches(0), Inches(3.5), PURPLE, 0.04)
    heading(s, "KẾT LUẬN")

    stat(s, Inches(0.3), Inches(1.25), Inches(2.2), Inches(1.0), "78.6%", "Bandwidth Saved", MINT)
    stat(s, Inches(2.7), Inches(1.25), Inches(2.2), Inches(1.0), "1.14 MB", "BF Size (FPR=1%)", PURPLE_L)
    stat(s, Inches(5.1), Inches(1.25), Inches(2.2), Inches(1.0), "983×", "Leverage", GOLD)
    stat(s, Inches(7.5), Inches(1.25), Inches(2.2), Inches(1.0), "100%", "Correctness", MINT)

    card(s, Inches(0.3), Inches(2.6), Inches(4.5), Inches(2.4),
         "🏆  Đóng góp của đồ án",
         ["", "✓  Bloom Filter cài từ scratch (from scratch)",
          "✓  Pipeline Semi-Join 5 bước đầy đủ",
          "✓  Phân tích lý thuyết + thực nghiệm",
          "✓  Web Dashboard tương tác",
          "✓  Demo Fault Tolerance (Kill Node B)",
          "✓  Liên kết chặt với Özsu & Valduriez"],
         ac=MINT, bsz=12)

    card(s, Inches(5.2), Inches(2.6), Inches(4.5), Inches(2.4),
         "🚀  Hạn chế & Hướng phát triển",
         ["", "•  Demo 100K/1M → Streaming cho 1M/10M",
          "•  Mô phỏng 1 máy → Docker / gRPC thực",
          "•  BF cơ bản → Counting / Cuckoo Filter",
          "•  Chỉ equi-join → Range join, multi-attr",
          "•  Chưa có replication → High availability"],
         ac=CORAL, bsz=12)

    mtxt(s, Inches(0.3), Inches(5.3), Inches(9.4), Inches(1.5), [
        {'text': '📚  Tài liệu tham khảo', 'sz': 12, 'c': WHITE, 'bold': True, 'spa': 4},
        {'text': '[1]  Özsu & Valduriez (2020). Principles of Distributed Database Systems, 4th Ed. — Chương 7.', 'sz': 10, 'c': DIM, 'spa': 2},
        {'text': '[2]  Bloom, B.H. (1970). Space/time trade-offs in hash coding with allowable errors. CACM.', 'sz': 10, 'c': DIM, 'spa': 2},
        {'text': '[3]  Kirsch & Mitzenmacher (2006). Less Hashing, Same Performance. ESA.', 'sz': 10, 'c': DIM},
    ], ls=1.1)
    page_num(s, 13, T)

    # ──────────── SLIDE 14: THANK YOU ────────────
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(s)
    deco_circle(s, Inches(6.5), Inches(-1), Inches(5), PURPLE, 0.06)
    deco_circle(s, Inches(-1), Inches(4), Inches(4), MINT, 0.04)
    deco_circle(s, Inches(7), Inches(5), Inches(3), GOLD, 0.03)
    accent_bar(s, Inches(0), Inches(0), Inches(10), PURPLE, Pt(4))
    accent_bar(s, Inches(0), Inches(7.5) - Pt(4), Inches(10), PURPLE, Pt(4))

    mtxt(s, Inches(0.5), Inches(1.5), Inches(9), Inches(5), [
        {'text': 'CẢM ƠN', 'sz': 52, 'c': WHITE, 'bold': True, 'fn': 'Segoe UI Semibold', 'align': PP_ALIGN.CENTER, 'spa': 0},
        {'text': 'THẦY / CÔ', 'sz': 52, 'c': PURPLE_L, 'bold': True, 'fn': 'Segoe UI Semibold', 'align': PP_ALIGN.CENTER, 'spa': 20},
        {'text': 'ĐÃ LẮNG NGHE', 'sz': 28, 'c': GRAY, 'align': PP_ALIGN.CENTER, 'spa': 30},
        {'text': 'Bloom Filter Semi-Join Optimizer', 'sz': 17, 'c': GOLD, 'bold': True, 'align': PP_ALIGN.CENTER, 'spa': 6},
        {'text': 'Đề tài #15  ·  Cơ Sở Dữ Liệu Phân Tán', 'sz': 13, 'c': DIM, 'align': PP_ALIGN.CENTER, 'spa': 30},
        {'text': '❓  Q & A  —  Xin mời thầy/cô đặt câu hỏi', 'sz': 16, 'c': MINT_L, 'bold': True, 'align': PP_ALIGN.CENTER},
    ], ls=1.1)
    page_num(s, 14, T)

    # ──────────── SAVE ────────────
    fp = os.path.join(OUTPUT_DIR, "Slide_Trinh_Bay_v2.pptx")
    prs.save(fp)
    print(f"\n  [OK] Da tao: {fp}")
    print(f"  Tong so slide: {T}")
    return fp


if __name__ == '__main__':
    print("\n" + "=" * 60)
    print("  TAO SLIDE POWERPOINT v2 — PREMIUM DESIGN")
    print("=" * 60)
    build()
    print("\n" + "=" * 60)
    print("  HOAN THANH! Mo file Slide_Trinh_Bay.pptx de xem.")
    print("=" * 60 + "\n")
